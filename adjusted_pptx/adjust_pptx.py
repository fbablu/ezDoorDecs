import os
import pandas as pd
import requests
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN


def adjust_pptx(residents_df, image_urls):
    prs = Presentation()
    folder_path = 'adjusted_pptx'
    image_dir = 'villager_images'
    os.makedirs(folder_path, exist_ok=True)
    os.makedirs(image_dir, exist_ok=True)

    paths = []

    for slide_index in range(0, len(residents_df) // 3 + 1):
        print(f"Processing slide {slide_index + 1}")

        slide_layout = prs.slide_layouts[5]  # Choosing a blank slide
        slide = prs.slides.add_slide(slide_layout)

        # Add three vertical rectangles inside the slide
        for i in range(3):
            resident_index = slide_index * 3 + i
            if resident_index < len(residents_df):
                resident = residents_df.iloc[resident_index]
                name = resident['Name']
                room = resident['Room']

                # Adjusted positioning and sizing
                left = Inches(0.15 + i * 3.3)
                top = Inches(0.15)
                width = Inches(3.04)
                height = Inches(7.1)

                # Add shape to the slide
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    left, top, width, height
                )

                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(
                    200, 200, 200)  # Light gray fill
                shape.line.color.rgb = RGBColor(0, 0, 0)  # Black border
                shape.line.width = Pt(4)

                # Add image
                try:
                    # Get image URL
                    image_url = image_urls.iloc[resident_index]
                    response = requests.get(image_url)
                    response.raise_for_status()  # Check for HTTP errors
                    image = Image.open(BytesIO(response.content))

                    # Convert RGBA to RGB if needed
                    if image.mode == 'RGBA':
                        image = image.convert('RGB')

                    # Save image temporarily
                    image_filename = os.path.join(
                        image_dir, f'temp_{name}.jpg')
                    image.save(image_filename)

                    img_border_left = Inches(left.inches + 0.275)
                    img_border_top = Inches(top.inches + 0.5)
                    img_border_width = Inches(2.5)
                    img_border_height = Inches(2.5)

                    img_border = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        img_border_left, img_border_top, img_border_width, img_border_height
                    )

                    img_border.fill.solid()
                    img_border.fill.fore_color.rgb = RGBColor(
                        200, 200, 200)  # Light gray fill
                    img_border.line.color.rgb = RGBColor(
                        0, 0, 0)  # Black border
                    img_border.line.width = Pt(4)

                    # Add villager image on top
                    slide.shapes.add_picture(image_filename,
                                             img_border_left,
                                             img_border_top,
                                             width=Inches(2.5),
                                             height=Inches(2.5))

                except Exception as e:
                    print(f"An error occurred with {name}'s image: {e}")

                # Add rectangle border for resident's name
                name_border_left = Inches(left.inches + 0.15)
                name_border_top = Inches(top.inches + 3.2)
                name_border_width = Inches(2.74)
                name_border_height = Inches(1)

                name_border = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    name_border_left, name_border_top, name_border_width, name_border_height
                )
                name_border.fill.solid()
                name_border.fill.fore_color.rgb = RGBColor(249, 245, 223)
                name_border.line.color.rgb = RGBColor(0, 0, 0)
                name_border.line.width = Pt(3)

                # Add text for resident's name on top of the border
                name_textbox = slide.shapes.add_textbox(
                    name_border_left, name_border_top, name_border_width, name_border_height)
                name_textframe = name_textbox.text_frame
                name_textframe.text = name
                name_textframe.paragraphs[0].alignment = PP_ALIGN.CENTER

                name_font = name_textframe.paragraphs[0].font
                name_font.name = 'Perpetua'
                name_font.size = Pt(66)
                name_font.bold = True
                name_font.italic = True
                name_font.color.rgb = RGBColor(0, 0, 0)

                # Add ellipse border for room number
                room_border_left = Inches(left.inches + 0.1)
                room_border_top = Inches(top.inches + 5)
                room_border_width = Inches(2.8)
                room_border_height = Inches(2.0)

                room_border = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    room_border_left, room_border_top, room_border_width, room_border_height
                )
                room_border.fill.solid()
                room_border.fill.fore_color.rgb = RGBColor(249, 245, 223)
                room_border.line.color.rgb = RGBColor(0, 0, 0)
                room_border.line.width = Pt(0)

                # Set the position of the room number text box
                room_textbox_left = Inches(left.inches + 0.2)
                room_textbox_top = Inches(5.64)

                # Create the text box with the updated positions
                room_textbox = slide.shapes.add_textbox(
                    room_textbox_left, room_textbox_top, room_border_width, room_border_height)
                room_textframe = room_textbox.text_frame

                # Add the room number text
                room_textframe.text = f"{room}"
                room_textframe.paragraphs[0].alignment = PP_ALIGN.CENTER

                # Set the font properties to match the resident name's font
                room_font = room_textframe.paragraphs[0].font
                room_font.name = 'Perpetua'
                room_font.size = Pt(66)
                room_font.bold = True
                room_font.italic = True
                room_font.color.rgb = RGBColor(0, 0, 0)

                # Adjust line spacing to move the text down if needed
                paragraph = room_textframe.paragraphs[0]
                paragraph.space_before = Pt(20)

                # Add bells icon
                bells_left = Inches(left.inches - 0.05)
                bells_top = Inches(top.inches + 3.9)
                bells_width = Inches(2.04)
                bells_height = Inches(2.04)

                bells_icon = os.path.join('bells.png')
                slide.shapes.add_picture(
                    bells_icon, bells_left, bells_top, bells_width, bells_height)

    # Save the modified presentation
    save_path = os.path.join(
        folder_path, 'MAIN_Adjusted_Residents_Presentation.pptx')
    prs.save(save_path)
    print(f"Presentation adjusted and saved as {save_path}")

    # Save the paths to a CSV file
    paths_df = pd.DataFrame(paths)
    paths_df.to_csv('image_paths.csv', index=False)
    print("Image paths saved to image_paths.csv")


# Load residents from CSV and image URLs from JSON
residents_df = pd.read_csv('residents_moore.csv')
image_urls = pd.read_json('villager_image_urls.json', typ='series')

# Call the function
adjust_pptx(residents_df, image_urls)
