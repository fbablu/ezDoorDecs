import os
import pandas as pd
import requests
from PIL import Image, ImageDraw
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
# Clash Royale card rarity colors
CARD_COLORS = {
    'common': {'primary': RGBColor(169, 169, 169), 'secondary': RGBColor(211, 211, 211)},
    'rare': {'primary': RGBColor(255, 140, 0), 'secondary': RGBColor(255, 165, 0)},
    'epic': {'primary': RGBColor(128, 0, 128), 'secondary': RGBColor(153, 50, 204)},
    'legendary': {'primary': RGBColor(255, 215, 0), 'secondary': RGBColor(255, 255, 0)},
    'champion': {'primary': RGBColor(255, 255, 0), 'secondary': RGBColor(255, 215, 0)}
}


def extract_dominant_colors(image_path, num_colors=3):
    """Extract dominant colors from an image using simple sampling"""
    try:
        # Load and resize image for faster processing
        image = Image.open(image_path)
        image = image.resize((50, 50))

        # Convert to RGB if necessary
        if image.mode != 'RGB':
            image = image.convert('RGB')

        # Alternative approach: sample pixels directly without getdata()
        width, height = image.size
        unique_colors = []

        # Sample pixels from different positions
        for y in range(0, height, 5):  # Every 5th row
            for x in range(0, width, 5):  # Every 5th column
                try:
                    pixel = image.getpixel((x, y))
                    if isinstance(pixel, tuple) and len(pixel) >= 3:
                        # Check if this color is significantly different from existing ones
                        is_unique = True
                        for existing in unique_colors:
                            if abs(pixel[0] - existing[0]) < 50 and abs(pixel[1] - existing[1]) < 50 and abs(pixel[2] - existing[2]) < 50:
                                is_unique = False
                                break
                        if is_unique and len(unique_colors) < num_colors:
                            # Take only RGB values
                            unique_colors.append(pixel[:3])

                        if len(unique_colors) >= num_colors:
                            break
                except (IndexError, ValueError):
                    continue
            if len(unique_colors) >= num_colors:
                break

        # Fill with default colors if not enough unique colors found
        while len(unique_colors) < num_colors:
            unique_colors.append((100 + len(unique_colors) * 50, 150, 200))

        return [RGBColor(int(c[0]), int(c[1]), int(c[2])) for c in unique_colors[:num_colors]]
    except Exception as e:
        print(f"Error extracting colors: {e}")
        # Return default gradient colors
        return [RGBColor(100, 150, 200), RGBColor(150, 100, 200), RGBColor(200, 100, 150)]


def create_gradient_background(width, height, rgb_colors):
    """Create a gradient background image using RGB color tuples"""
    image = Image.new('RGB', (width, height))
    draw = ImageDraw.Draw(image)

    # Convert RGBColor objects to tuples
    colors = []
    for color in rgb_colors:
        # Extract RGB values from the hex representation
        hex_val = hex(color)[2:].zfill(6)
        r = int(hex_val[0:2], 16)
        g = int(hex_val[2:4], 16)
        b = int(hex_val[4:6], 16)
        colors.append((r, g, b))

    if len(colors) < 2:
        colors = [(100, 100, 100), (200, 200, 200)]

    for y in range(height):
        # Calculate which colors to blend and the ratio
        section = (y / height) * (len(colors) - 1)
        color1_idx = int(section)
        color2_idx = min(color1_idx + 1, len(colors) - 1)
        ratio = section - color1_idx

        color1 = colors[color1_idx]
        color2 = colors[color2_idx]

        # Interpolate between the two colors
        r = int(color1[0] * (1 - ratio) + color2[0] * ratio)
        g = int(color1[1] * (1 - ratio) + color2[1] * ratio)
        b = int(color1[2] * (1 - ratio) + color2[2] * ratio)

        draw.line([(0, y), (width, y)], fill=(r, g, b))

    return image


def fetch_and_save_image(image_url, image_filename):
    """Fetch and save card image"""
    try:
        response = requests.get(image_url)
        response.raise_for_status()
        image = Image.open(BytesIO(response.content))

        # Convert RGBA to RGB if needed
        if image.mode == 'RGBA':
            white_bg = Image.new('RGB', image.size, (255, 255, 255))
            white_bg.paste(image, mask=image.split()
                           [-1] if len(image.split()) == 4 else None)
            image = white_bg

        # Save the image locally
        image.save(image_filename)
        return image_filename
    except Exception as e:
        print(f"Error fetching image: {e}")
        return None


def create_clash_royale_presentation(residents_df, card_data):
    """Create Clash Royale themed presentation"""
    prs = Presentation()
    folder_path = 'clash_royale_images'
    os.makedirs(folder_path, exist_ok=True)

    for i in range(0, len(residents_df), 3):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide

        # Add arena-style background
        main_left = Inches(0.3)
        main_top = Inches(0.3)
        main_width = Inches(9.4)
        main_height = Inches(7)

        # Create arena background with blue gradient
        main_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            main_left, main_top, main_width, main_height
        )
        main_shape.fill.solid()
        main_shape.fill.fore_color.rgb = RGBColor(30, 144, 255)  # Arena blue
        main_shape.line.color.rgb = RGBColor(25, 25, 112)  # Dark blue border
        main_shape.line.width = Pt(4)

        # Add up to three residents per slide
        for j in range(3):
            index = i + j
            if index < len(residents_df):
                resident = residents_df.iloc[index]
                name = resident['Name']
                room = resident['Room']

                # Get card data for this resident
                card = card_data[index % len(card_data)] if index < len(
                    card_data) else card_data[0]
                rarity = card['rarity']
                colors = CARD_COLORS[rarity]

                left = Inches(0.8 + j * 3)
                top = Inches(0.8)
                width = Inches(2.6)
                height = Inches(5.8)

                # Download and save the card image first to extract colors
                card_image_filename = os.path.join(
                    folder_path, f'{name}_card.jpg')
                if fetch_and_save_image(card['image_url'], card_image_filename):
                    # Extract dominant colors from the card image
                    dominant_colors = extract_dominant_colors(
                        card_image_filename)

                    # Create gradient background image
                    gradient_bg = create_gradient_background(
                        int(width.inches * 100), int(height.inches * 100), dominant_colors)
                    gradient_filename = os.path.join(
                        folder_path, f'{name}_gradient.png')
                    gradient_bg.save(gradient_filename)

                    # Add gradient background to card
                    try:
                        slide.shapes.add_picture(
                            gradient_filename,
                            left, top, width, height
                        )
                    except Exception as e:
                        print(f"Error adding gradient background: {e}")

                # Create card-style shape with rarity colors
                card_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left, top, width, height
                )
                card_shape.fill.solid()
                card_shape.fill.fore_color.rgb = colors['primary']
                card_shape.line.color.rgb = RGBColor(
                    255, 255, 255)  # White border
                card_shape.line.width = Pt(3)

                # Add Elixir.png icon in top left corner
                try:
                    elixir_left = Inches(left.inches + 0.1)
                    elixir_top = Inches(top.inches + 0.1)
                    elixir_width = Inches(0.6)
                    elixir_height = Inches(0.6)

                    slide.shapes.add_picture(
                        'Elixir.png', elixir_left, elixir_top, elixir_width, elixir_height)
                except Exception as e:
                    print(f"Error adding Elixir.png: {e}")

                # Add inner card area
                inner_left = Inches(left.inches + 0.1)
                inner_top = Inches(top.inches + 0.8)
                inner_width = Inches(width.inches - 0.2)
                inner_height = Inches(2.2)

                inner_shape = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    inner_left, inner_top, inner_width, inner_height
                )
                inner_shape.fill.solid()
                inner_shape.fill.fore_color.rgb = RGBColor(
                    240, 240, 240)  # Light gray
                inner_shape.line.color.rgb = RGBColor(
                    200, 200, 200)  # Gray border
                inner_shape.line.width = Pt(2)

                # Add card image if available
                if card_image_filename:
                    try:
                        slide.shapes.add_picture(
                            card_image_filename,
                            Inches(inner_left.inches + 0.15),
                            Inches(inner_top.inches + 0.15),
                            width=Inches(inner_width.inches - 0.3),
                            height=Inches(inner_height.inches - 0.3),
                        )
                    except Exception as e:
                        print(f"Error adding image for {name}: {e}")

                # Add name with Clash Royale style font
                name_top = Inches(top.inches + 3.2)
                name_textbox = slide.shapes.add_textbox(
                    Inches(
                        left.inches + 0.1), name_top, Inches(width.inches - 0.2), Inches(0.8)
                )
                name_textframe = name_textbox.text_frame
                name_textframe.text = name
                name_textframe.paragraphs[0].alignment = PP_ALIGN.CENTER

                # Clash Royale style font (bold, prominent)
                name_font = name_textframe.paragraphs[0].font
                name_font.name = 'Impact'  # Bold, game-like font
                name_font.size = Pt(18)
                name_font.bold = True
                name_font.color.rgb = RGBColor(255, 255, 255)  # White text

                # Add room number text (simple text display)
                room_textbox = slide.shapes.add_textbox(
                    Inches(left.inches + 0.1), Inches(top.inches + 4.2),
                    Inches(width.inches - 0.2), Inches(0.6)
                )
                room_textframe = room_textbox.text_frame
                room_textframe.text = f"{room}"
                room_textframe.paragraphs[0].alignment = PP_ALIGN.CENTER

                room_font = room_textframe.paragraphs[0].font
                room_font.name = 'Impact'
                room_font.size = Pt(16)
                room_font.bold = True
                room_font.color.rgb = RGBColor(255, 255, 255)  # White text

                # Add rarity indicator
                rarity_textbox = slide.shapes.add_textbox(
                    Inches(left.inches + 0.1), Inches(top.inches + 5.2),
                    Inches(width.inches - 0.2), Inches(0.4)
                )

                rarity_textframe = rarity_textbox.text_frame
                rarity_textframe.text = rarity.upper()
                rarity_textframe.paragraphs[0].alignment = PP_ALIGN.CENTER

                rarity_font = rarity_textframe.paragraphs[0].font
                rarity_font.name = 'Impact'
                rarity_font.size = Pt(10)
                rarity_font.bold = True
                rarity_font.color.rgb = colors['secondary']

    # Save the presentation
    save_path = 'Clash_Royale_Door_Decks.pptx'
    prs.save(save_path)
    print(f"Clash Royale presentation created: {save_path}")

    # Save image paths to CSV
    paths_df = pd.DataFrame({
        'Name': [resident['Name'] for _, resident in residents_df.iterrows()],
        'Path': [os.path.join(folder_path, f'{resident["Name"]}_card.jpg') for _, resident in residents_df.iterrows()],
        'Rarity': [card_data[i % len(card_data)]['rarity'] for i in range(len(residents_df))]
    })
    paths_df.to_csv('clash_royale_image_paths.csv', index=False)
    print("Image paths saved to clash_royale_image_paths.csv")


if __name__ == "__main__":
    # Load residents data
    residents_df = pd.read_csv('residents_moore.csv')

    try:
        with open('clash_royale_card_data.json', 'r') as f:
            import json
            card_data = json.load(f)
    except FileNotFoundError:
        print("Please run get_clash_royale_images.py first to generate card data")
        card_data = []

    # Create the presentation
    create_clash_royale_presentation(residents_df, card_data)
